from fastapi import APIRouter, HTTPException, Depends, UploadFile, File, Form
import httpx
import json
import io
import os
import re
from pydantic import BaseModel
from utils.auth import admin_required
from docx import Document
from pptx import Presentation
from config.db import get_database

router = APIRouter(prefix="/api/ai", tags=["AI Generation"])
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434/api/generate")
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "llama3")
DEFAULT_AI_SETTINGS = {
    "enabled": True,
    "default_prompt": "Generate concise and accurate multiple-choice questions.",
    "difficulty": "medium",
    "questions_limit": 20,
}

class GenerateRequest(BaseModel):
    topic: str
    num_questions: int


async def get_ai_settings(db=None) -> dict:
    db = db or get_database()
    if not db:
        return DEFAULT_AI_SETTINGS.copy()
    settings_col = db["quizzify"]["ai_settings"]
    doc = await settings_col.find_one({"_id": "default"})
    if not doc:
        await settings_col.update_one(
            {"_id": "default"},
            {"$set": DEFAULT_AI_SETTINGS},
            upsert=True,
        )
        return DEFAULT_AI_SETTINGS.copy()

    return {
        "enabled": bool(doc.get("enabled", DEFAULT_AI_SETTINGS["enabled"])),
        "default_prompt": str(doc.get("default_prompt", DEFAULT_AI_SETTINGS["default_prompt"])).strip(),
        "difficulty": str(doc.get("difficulty", DEFAULT_AI_SETTINGS["difficulty"])).strip().lower() or "medium",
        "questions_limit": int(doc.get("questions_limit", DEFAULT_AI_SETTINGS["questions_limit"])),
    }


def _normalize_questions(parsed_questions: list, max_questions: int) -> list:
    normalized = []
    for item in parsed_questions:
        if not isinstance(item, dict):
            continue

        text = str(item.get("text", "")).strip()
        options = item.get("options", [])
        correct_answer = str(item.get("correct_answer", "")).strip()

        if not text or not isinstance(options, list):
            continue

        cleaned_options = [str(opt).strip() for opt in options if str(opt).strip()]
        unique_options = []
        for opt in cleaned_options:
            if opt not in unique_options:
                unique_options.append(opt)

        if correct_answer and correct_answer not in unique_options:
            unique_options.insert(0, correct_answer)

        while len(unique_options) < 4:
            unique_options.append(f"None of the above ({len(unique_options) + 1})")

        unique_options = unique_options[:4]
        if not correct_answer or correct_answer not in unique_options:
            correct_answer = unique_options[0]

        normalized.append({
            "text": text,
            "options": unique_options,
            "correct_answer": correct_answer,
        })

        if len(normalized) >= max_questions:
            break

    return normalized


async def _generate_with_ollama(prompt: str, num_questions: int) -> list:
    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(
            OLLAMA_URL,
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "format": "json",
            },
        )
        response.raise_for_status()
        data = response.json()
        response_text = data.get("response", "")

        parsed_questions = json.loads(response_text)
        if isinstance(parsed_questions, dict) and "questions" in parsed_questions:
            parsed_questions = parsed_questions["questions"]

        if not isinstance(parsed_questions, list):
            raise ValueError("Expected a JSON array")

        normalized = _normalize_questions(parsed_questions, max_questions=num_questions)
        if not normalized:
            raise ValueError("No valid questions returned by model")
        return normalized


def _build_topic_prompt(topic: str, num_questions: int, settings: dict) -> str:
    return f"""{settings.get("default_prompt", DEFAULT_AI_SETTINGS["default_prompt"])}
Difficulty level: {settings.get("difficulty", "medium")}.
Generate {num_questions} multiple-choice questions about the topic: "{topic}".

Return ONLY a valid JSON array of objects, with no markdown, no backticks, and no extra text.
Each object must have the following exact keys:
- "text": The question text as a string
- "options": An array of exactly 4 strings (one correct, three plausible distractors)
- "correct_answer": A string matching exactly one of the options
"""


def _build_text_prompt(text: str, num_questions: int, settings: dict) -> str:
    return f"""{settings.get("default_prompt", DEFAULT_AI_SETTINGS["default_prompt"])}
Difficulty level: {settings.get("difficulty", "medium")}.
Generate {num_questions} multiple-choice questions based ONLY on the following content:
"{text}"

Return ONLY a valid JSON array of objects, with no markdown, no backticks, and no extra text.
Each object must have the following exact keys:
- "text": The question text as a string
- "options": An array of exactly 4 strings (one correct, three plausible distractors)
- "correct_answer": A string matching exactly one of the options
"""


async def extract_text_from_upload(file: UploadFile) -> str:
    text = ""
    filename = (file.filename or "").lower()
    content = await file.read()

    if filename.endswith(".pdf"):
        try:
            import fitz
        except Exception:
            raise HTTPException(status_code=500, detail="PyMuPDF is required for PDF extraction")

        pdf = fitz.open(stream=content, filetype="pdf")
        for page in pdf:
            page_text = page.get_text("text")
            if page_text:
                text += page_text + "\n"
    elif filename.endswith(".docx"):
        doc = Document(io.BytesIO(content))
        for para in doc.paragraphs:
            if para.text:
                text += para.text + "\n"
    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
        prs = Presentation(io.BytesIO(content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    else:
        raise HTTPException(status_code=400, detail="Unsupported file format. Please upload PDF, DOCX, or PPTX.")

    text = text[:10000]
    if not text.strip():
        raise HTTPException(status_code=400, detail="Could not extract readable text from the uploaded file.")

    return text


async def generate_questions_from_text(text: str, num_questions: int, settings: dict) -> dict:
    prompt = _build_text_prompt(text=text, num_questions=num_questions, settings=settings)
    try:
        questions = await _generate_with_ollama(prompt, num_questions=num_questions)
        return {"questions": questions, "source": "ollama"}
    except Exception:
        questions = _fallback_questions_from_text(text, num_questions=num_questions)
        return {"questions": questions, "source": "fallback"}


def _sentences_from_text(text: str) -> list:
    raw = re.split(r"(?<=[.!?])\s+", text)
    sentences = []
    seen = set()
    for part in raw:
        cleaned = " ".join(part.split()).strip()
        if len(cleaned) < 25:
            continue
        key = cleaned.lower()
        if key in seen:
            continue
        seen.add(key)
        sentences.append(cleaned)
    return sentences


def _fallback_questions_from_topic(topic: str, num_questions: int) -> list:
    topic = topic.strip() or "this topic"
    templates = [
        {
            "text": f"Which statement best describes {topic}?",
            "correct_answer": f"{topic} includes foundational concepts that can be applied in real scenarios.",
            "options": [
                f"{topic} includes foundational concepts that can be applied in real scenarios.",
                f"{topic} has no practical applications and is only theoretical.",
                f"{topic} is unrelated to problem solving.",
                f"{topic} cannot be learned through structured practice.",
            ],
        },
        {
            "text": f"What is a recommended way to improve understanding of {topic}?",
            "correct_answer": f"Study core ideas of {topic} and apply them through examples.",
            "options": [
                f"Study core ideas of {topic} and apply them through examples.",
                f"Avoid practicing {topic} to prevent confusion.",
                f"Memorize random facts without connecting them to {topic}.",
                f"Ignore feedback when learning {topic}.",
            ],
        },
        {
            "text": f"Which option is most likely true about mastery in {topic}?",
            "correct_answer": f"Mastery in {topic} improves with consistent practice and revision.",
            "options": [
                f"Mastery in {topic} improves with consistent practice and revision.",
                f"Mastery in {topic} happens instantly without effort.",
                f"Mastery in {topic} is impossible for beginners.",
                f"Mastery in {topic} depends only on luck.",
            ],
        },
    ]

    questions = []
    for i in range(num_questions):
        questions.append(templates[i % len(templates)])
    return questions


def _fallback_questions_from_text(text: str, num_questions: int) -> list:
    sentences = _sentences_from_text(text)
    if not sentences:
        return _fallback_questions_from_topic("the provided material", num_questions)

    questions = []
    for i in range(min(num_questions, len(sentences))):
        correct = sentences[i]
        distractors = []
        for candidate in sentences:
            if candidate != correct and candidate not in distractors:
                distractors.append(candidate)
            if len(distractors) == 3:
                break

        while len(distractors) < 3:
            distractors.append(f"This statement is not supported by the provided material ({len(distractors) + 1}).")

        options = [correct, distractors[0], distractors[1], distractors[2]]
        questions.append(
            {
                "text": "Based on the uploaded content, which statement is correct?",
                "options": options,
                "correct_answer": correct,
            }
        )

    while len(questions) < num_questions:
        questions.extend(_fallback_questions_from_topic("the provided material", 1))

    return questions[:num_questions]

@router.post("/generate-quiz")
async def generate_quiz(req: GenerateRequest, current_user: dict = Depends(admin_required)):
    db = get_database()
    settings = await get_ai_settings(db)
    if not settings.get("enabled", True):
        raise HTTPException(status_code=403, detail="AI generation is currently disabled by admin settings")

    num_questions = max(1, min(int(req.num_questions), int(settings.get("questions_limit", 20))))
    prompt = _build_topic_prompt(topic=req.topic, num_questions=num_questions, settings=settings)
    try:
        questions = await _generate_with_ollama(prompt, num_questions=num_questions)
        return {"questions": questions, "source": "ollama"}
    except Exception:
        # If local LLM is unavailable, return deterministic fallback questions.
        questions = _fallback_questions_from_topic(req.topic, num_questions=num_questions)
        return {"questions": questions, "source": "fallback"}

@router.post("/generate-from-file")
async def generate_from_file(
    file: UploadFile = File(...),
    num_questions: int = Form(...),
    current_user: dict = Depends(admin_required)
):
    db = get_database()
    settings = await get_ai_settings(db)
    if not settings.get("enabled", True):
        raise HTTPException(status_code=403, detail="AI generation is currently disabled by admin settings")

    num_questions = max(1, min(int(num_questions), int(settings.get("questions_limit", 20))))
    
    try:
        text = await extract_text_from_upload(file)
        return await generate_questions_from_text(text=text, num_questions=num_questions, settings=settings)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate quiz from file: {str(e)}")
